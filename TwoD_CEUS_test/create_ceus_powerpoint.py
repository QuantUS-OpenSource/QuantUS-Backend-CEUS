"""
Script to create a PowerPoint presentation with CEUS frames and motion-corrected ROI overlay.
Each slide contains one frame with the ROI overlaid.
"""

import numpy as np
import nibabel as nib
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import matplotlib.pyplot as plt
import io
from pathlib import Path
from tqdm import tqdm


def load_and_process_roi(roi_path, target_shape):
    """Load and process ROI mask to match target shape."""
    roi_seg = nib.load(roi_path).get_fdata()

    # Process combined side-by-side visualization if needed
    if roi_seg.shape[0] == 1048:
        print(f"  Detected combined side-by-side ROI, splitting...")
        mid = 524
        left_half = roi_seg[:mid, :, :]
        right_half = roi_seg[mid:, :, :]

        # Choose half with more ROI data
        if np.sum(left_half > 0) > np.sum(right_half > 0):
            roi_seg = left_half
        else:
            roi_seg = right_half

        # Transpose to match image orientation
        if roi_seg.shape != target_shape:
            roi_seg = roi_seg.transpose(1, 0, 2)

    return roi_seg


def normalize_ceus_data(ceus_data):
    """Normalize CEUS data for visualization."""
    print("Normalizing CEUS data...")
    vmin = np.percentile(ceus_data, 0.1)
    vmax = np.percentile(ceus_data, 99.9)
    normalized = np.clip(ceus_data, vmin, vmax)
    normalized = ((normalized - vmin) / (vmax - vmin) * 255).astype(np.uint8)
    return normalized


def create_frame_with_overlay(ceus_frame, roi_mask, frame_idx, alpha=0.5):
    """Create a matplotlib figure with CEUS frame and ROI overlay."""
    fig, ax = plt.subplots(figsize=(10, 10), dpi=100)

    # Show CEUS frame
    ax.imshow(ceus_frame, cmap='gray', vmin=0, vmax=255)

    # Overlay ROI if present
    if np.any(roi_mask > 0):
        ax.imshow(roi_mask, alpha=alpha, cmap='Reds', vmin=0, vmax=1)

    ax.set_title(f'Frame {frame_idx}', fontsize=16, fontweight='bold')
    ax.axis('off')
    plt.tight_layout()

    # Convert to PIL Image
    buf = io.BytesIO()
    plt.savefig(buf, format='png', bbox_inches='tight', dpi=100)
    buf.seek(0)
    img = Image.open(buf)
    plt.close(fig)

    return img


def create_powerpoint_with_frames(
    ceus_path,
    roi_path,
    output_path,
    start_frame=0,
    end_frame=None,
    frame_step=1,
    alpha=0.5,
    slide_width_inches=10,
    slide_height_inches=7.5
):
    """
    Create a PowerPoint presentation with CEUS frames and ROI overlay.

    Parameters
    ----------
    ceus_path : str
        Path to CEUS NIfTI file
    roi_path : str
        Path to motion-corrected ROI NIfTI file
    output_path : str
        Path for output PowerPoint file
    start_frame : int
        Starting frame index (default: 0)
    end_frame : int
        Ending frame index (default: None = last frame)
    frame_step : int
        Step between frames (default: 1 = every frame)
    alpha : float
        ROI overlay transparency (0-1, default: 0.5)
    slide_width_inches : float
        Slide width in inches (default: 10)
    slide_height_inches : float
        Slide height in inches (default: 7.5)
    """

    print("="*70)
    print("CREATING POWERPOINT WITH CEUS FRAMES AND ROI OVERLAY")
    print("="*70)

    # Load CEUS data
    print(f"\nLoading CEUS data from: {ceus_path}")
    ceus_data = nib.load(ceus_path).get_fdata()
    print(f"  Shape: {ceus_data.shape}")
    print(f"  Data type: {ceus_data.dtype}")

    # Load and process ROI
    print(f"\nLoading ROI from: {roi_path}")
    roi_seg = load_and_process_roi(roi_path, ceus_data.shape)
    print(f"  ROI shape: {roi_seg.shape}")
    print(f"  ROI pixels: {np.sum(roi_seg > 0)}")

    # Normalize CEUS data
    ceus_normalized = normalize_ceus_data(ceus_data)

    # Determine frame range
    total_frames = ceus_data.shape[2]
    if end_frame is None:
        end_frame = total_frames - 1

    frame_indices = range(start_frame, end_frame + 1, frame_step)
    num_slides = len(frame_indices)

    print(f"\nGenerating {num_slides} slides...")
    print(f"  Frame range: {start_frame} to {end_frame} (step: {frame_step})")

    # Create PowerPoint presentation
    prs = Presentation()
    prs.slide_width = Inches(slide_width_inches)
    prs.slide_height = Inches(slide_height_inches)

    # Add blank slide layout
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout

    # Generate slides
    for frame_idx in tqdm(frame_indices, desc="Creating slides"):
        # Get CEUS frame and corresponding ROI mask
        ceus_frame = ceus_normalized[:, :, frame_idx]

        # Get ROI mask for this frame
        if roi_seg.ndim == 3:
            # ROI has temporal dimension
            if roi_seg.shape[2] == ceus_data.shape[2]:
                # Same number of frames, use corresponding frame
                roi_mask = roi_seg[:, :, frame_idx]
            elif frame_idx < roi_seg.shape[2]:
                # ROI has fewer frames, use frame if available
                roi_mask = roi_seg[:, :, frame_idx]
            else:
                # Frame index beyond ROI frames, use last ROI frame
                roi_mask = roi_seg[:, :, -1]
        else:
            # ROI is 2D, use for all frames
            roi_mask = roi_seg

        # Create image with overlay
        img = create_frame_with_overlay(ceus_frame, roi_mask, frame_idx, alpha=alpha)

        # Save to temporary buffer
        img_buf = io.BytesIO()
        img.save(img_buf, format='PNG')
        img_buf.seek(0)

        # Add slide
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add image to slide (centered)
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(slide_width_inches - 1)
        height = Inches(slide_height_inches - 1)

        pic = slide.shapes.add_picture(img_buf, left, top, width=width, height=height)

    # Save PowerPoint
    print(f"\nSaving PowerPoint to: {output_path}")
    prs.save(output_path)

    print(f"\n✓ PowerPoint created successfully!")
    print(f"  Total slides: {num_slides}")
    print(f"  File: {output_path}")
    print(f"  File size: {Path(output_path).stat().st_size / (1024*1024):.2f} MB")
    print("="*70)


if __name__ == "__main__":
    # Example usage - update these paths for your data

    # Path to CEUS NIfTI file
    ceus_path = '/Users/samantha/Desktop/ultrasound lab stuff/raw_ctdna/p14/wk12/combined_nifti/p14_wk12_CHI_RAW.nii'

    # Path to motion-corrected ROI NIfTI file
    roi_path = '/Users/samantha/Desktop/ultrasound lab stuff/raw_mc_ctdna/p14/wk12/p14_wk12_raw_mc_roi.nii.gz'

    # Output PowerPoint file
    output_path = '/Users/samantha/Desktop/ceus_frames_with_roi.pptx'

    # Create PowerPoint
    # You can adjust these parameters:
    # - start_frame: which frame to start at
    # - end_frame: which frame to end at (None = last frame)
    # - frame_step: include every Nth frame (1 = all frames, 5 = every 5th frame)
    # - alpha: ROI transparency (0 = invisible, 1 = opaque)

    create_powerpoint_with_frames(
        ceus_path='/Users/samantha/Desktop/ultrasound lab stuff/raw_ctdna/p14/wk12/combined_nifti/p14_wk12_CHI_RAW.nii',
        roi_path='/Users/samantha/Desktop/ultrasound lab stuff/raw_mc_ctdna/p14/wk12/p14_wk12_raw_mc_roi.nii.gz',
        output_path='/Users/samantha/Desktop/ultrasound lab stuff/p14_wk12_frames.pptx',
        start_frame=0,      # Start at frame 0
        end_frame=500,     # Go to last frame (or specify a number)
        frame_step=1,       # Include every frame (change to 5 for every 5th frame)
        alpha=0.5,          # 50% transparency for ROI overlay
    )
